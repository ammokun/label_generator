 #-*- coding: utf-8 -*-

from pptx import Presentation
from pptx.util import Inches
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.shapes.connector import Connector
from pptx.dml.line import LineFormat
import openpyxl as pyxl
import os
import datetime

count = 0
#os.chdir('')
os.getcwd()

#産地コードごとにまとめて出力する機能
santi_list=["AAA沢","BBB沢","CCC沢"]

#1ページに書き込むラベルの数 rowは2 or3　columnは6以下を推奨　文字の大きさによって変更
label_row=2
label_column=6
label_num=label_row*label_column

x_zahyou=[]
y_zahyou=[]
for i in range(0,label_num):
    x_zahyou.append((i%label_row)*1/label_row)
    y_zahyou.append(int(i/label_row)*1/label_column)

temp=[]

#エクセルのファイル名
wb=pyxl.load_workbook('sample.xlsx')
sheet=wb['Sheet1']

hyohon_list=[]

print('Reading...\n')
for row in range(1,sheet.max_row+1):
	

	hyohon={

	'scientific_name_genes':sheet['A'+str(row)].value,
	'scientific_name_species':sheet['B'+str(row)].value,
	'scientific_name_roman':sheet['C'+str(row)].value,
    'japanese_name':sheet['D'+str(row)].value,
	'loc':sheet['E'+str(row)].value,
	'stratum':sheet['F'+str(row)].value,
	'age':sheet['G'+str(row)].value,
	'date':sheet['H'+str(row)].value,
	'col_name':sheet['I'+str(row)].value,
	'size':sheet['J'+str(row)].value,
    'remarks':sheet['K'+str(row)].value,
	'file_name':sheet['L'+str(row)].value,
    'number':sheet['M'+str(row)].value,
    'loc_code':sheet['N'+str(row)].value,

	}

	hyohon_list.append(hyohon)
max=len(hyohon_list)
print('Excel file has been successfully read! max=',max)
prs = Presentation()
blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)

prs.slide_width=7560000
prs.slide_height=10656000

width=prs.slide_width
height=prs.slide_height

def if_nonetype(table_name):
    if isinstance(table_name,type(None)) == True:
        return str(' ')
    else:
        return table_name
def get_or_add_ln(self):
    return self._element.spPr.get_or_add_ln()

Connector.get_or_add_ln = get_or_add_ln

def line(x1,y1,x2,y2,line_width):
    shapes=slide.shapes
    connector = shapes.add_connector(
    MSO_CONNECTOR.STRAIGHT, int(width*x1), int(height*y1), int(width*x2), int(height*y2)
    )
    if not hasattr(connector, "ln"):
        connector.ln = connector.get_or_add_ln()
    line = LineFormat(connector)
    line.width=Pt(line_width)
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(0, 0, 0)

#切り取り線
def cut_box(x,y):
    line(x,y,x+1/label_row,y,1)
    line(x,y+1/label_column,x+1/label_row,y+1/label_column,1)
    line(x,y,x,y+1/label_column,1)
    line(x+1/label_row,y,x+1/label_row,y+1/label_column,1)

#ラベルの外枠
def label_box(x,y):
    line(x+0.1/label_row,y+0.1/label_column,x+0.9/label_row,y+0.1/label_column,2)
    line(x+0.1/label_row,y+0.9/label_column,x+0.9/label_row,y+0.9/label_column,2)
    line(x+0.1/label_row,y+0.1/label_column,x+0.1/label_row,y+0.9/label_column,2)
    line(x+0.9/label_row,y+0.1/label_column,x+0.9/label_row,y+0.9/label_column,2)

#ラベルの内容
def add_lavel(x,y,number):


    date1=hyohon_list[number]['date']
    col_date=date1.strftime('%Y/%m/%d')

    txBox = slide.shapes.add_textbox(0, 0, Inches(1), Inches(1))
    txBox.left=int(width*x+width*0.05)
    #txBox.top=int((height*0.95)*y+Inches(2.1))

    txBox.top=int(height*y)
    tf = txBox.text_frame
    tf.clear()
    #tf.text=hyohon_list[number]['japanese_name']
    
    if hyohon_list[number]['scientific_name_genes'] is None:
        p = tf.add_paragraph()
        
        roman = p.add_run()
        roman.text = hyohon_list[number]['scientific_name_roman']+" "
        #p.font.bold = True
        roman.font.roman = True
        roman.font.bold = True
        roman.font.name = "Times New Roman"

    elif hyohon_list[number]['scientific_name_roman'] is None:
        p = tf.add_paragraph()

        genes = p.add_run()
        genes.text = hyohon_list[number]['scientific_name_genes']+ " " 
        #p.font.bold = True
        genes.font.italic = True
        genes.font.bold = True
        genes.font.name = "Times New Roman"

        species = p.add_run()
        species.text = hyohon_list[number]['scientific_name_species']
        #p.font.bold = True
        species.font.italic = True
        species.font.bold = True
        species.font.name = "Times New Roman"

    elif hyohon_list[number]['scientific_name_species'] is None: 
        p = tf.add_paragraph()

        genes = p.add_run()
        genes.text = hyohon_list[number]['scientific_name_genes']+ " "
        #p.font.bold = True
        genes.font.italic = True
        genes.font.bold = True
        genes.font.name = "Times New Roman"
        roman = p.add_run()
        roman.text = hyohon_list[number]['scientific_name_roman']
        #p.font.bold = True
        roman.font.roman = True
        roman.font.bold = True
        roman.font.name = "Times New Roman"

    else :
        p = tf.add_paragraph()
        genes = p.add_run()
        genes.text = hyohon_list[number]['scientific_name_genes'] +" "
        #p.font.bold = True
        genes.font.italic = True
        genes.font.bold = True
        genes.font.name = "Times New Roman"

        roman = p.add_run()
        roman.text = hyohon_list[number]['scientific_name_roman']+" "
        #p.font.bold = True
        roman.font.roman = True
        roman.font.bold = True
        roman.font.name = "Times New Roman"

        species = p.add_run()
        species.text = hyohon_list[number]['scientific_name_species']
        #p.font.bold = True
        species.font.italic = True
        species.font.bold = True
        species.font.name = "Times New Roman"

    p = tf.add_paragraph()
    p.text = if_nonetype(hyohon_list[number]['japanese_name'])
    p.font.size=Pt(10.5)

    p = tf.add_paragraph()
    #p.text =hyohon_list[number]['col_name']#+'  '+hyohon_list[number]['size']
    p.text =hyohon_list[number]['loc']
    p.font.size=Pt(10.5)

    p = tf.add_paragraph()
    #p.text =hyohon_list[number]['col_name']#+'  '+hyohon_list[number]['size']
    p.text =hyohon_list[number]['age']
    p.font.size=Pt(10.5)

    p = tf.add_paragraph()
    #p.text =hyohon_list[number]['col_name']#+'  '+hyohon_list[number]['size']
    p.text =hyohon_list[number]['stratum']
    p.font.size=Pt(10.5)

    p = tf.add_paragraph()
    #p.text =hyohon_list[number]['col_name']#+'  '+hyohon_list[number]['size']
    p.text =str(col_date)+'  '+hyohon_list[number]['col_name']+'  '+hyohon_list[number]['size']
    p.font.size=Pt(10.5)

    p = tf.add_paragraph()
    #p.text =hyohon_list[number]['col_name']#+'  '+hyohon_list[number]['size']
    p.text =hyohon_list[number]['remarks']
    p.font.size=Pt(10.5)

#テキストの配置
def add_text(x,y,text,font_size,bold,italic):
    txBox = slide.shapes.add_textbox(0, 0, Inches(1), Inches(1))
    txBox.left = x
    txBox.top = y
    tf = txBox.text_frame
    tf.clear()
    p = tf.add_paragraph()
    p.text = text

    p.font.size=Pt(font_size)
    p.font.bold = bold
    p.font.italic = italic

#ラベルの位置決定と配置
for santi in santi_list:
    #santi = unicode(santi,'utf-8')
    del temp[:]
    for i in range(1,max):
        if hyohon_list[i]['loc_code'] == santi:
            temp.append(i)
    n=len(temp)
    print(santi)
    print(n)
    #print(temp)


    for i in temp:
        x=x_zahyou[(count%label_num)]
        y=y_zahyou[(count%label_num)]
    
        add_lavel(x,y,i)
        title_x = Inches(0.5)
        title_y = 0
       # title = santi+u'の化石'

        #add_text(title_x,title_y,title,28,True,False) #title
        #add_text(page_x,page_y,str(page),12,False,False)   #page number
    
        label_box(x,y)
        cut_box(x,y)


        count=count+1
        if count == n:
            slide = prs.slides.add_slide(blank_slide_layout)
            count = 0
            
            break
        if count%label_num == 0 :
            slide = prs.slides.add_slide(blank_slide_layout)
        
            
prs.save('test.pptx')

